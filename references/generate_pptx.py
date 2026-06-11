import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    DARK_BG = RGBColor(15, 23, 42)      # Deep Navy
    LIGHT_BG = RGBColor(248, 250, 252)  # Light Slate
    TEXT_LIGHT = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(100, 116, 139) # Slate Gray
    TEAL_COLOR = RGBColor(14, 165, 233)  # Sky Blue
    GREEN_COLOR = RGBColor(34, 197, 94)  # Emerald
    ACCENT_BOX = RGBColor(241, 245, 249) # Subtle Gray

    blank_layout = prs.slide_layouts[6]

    # -----------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -----------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    # Title Text Box
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "AZUREOPS SENTINEL"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEAL_COLOR
    p.font.name = "Segoe UI"
    p.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Enterprise GitOps Blueprint"
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = "Segoe UI"
    p2.space_after = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "A Multi-Source Event-Driven Declarative Remediation Architecture for Public Clouds"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Segoe UI"
    p3.space_after = Pt(40)

    p4 = tf.add_paragraph()
    p4.text = "Presented by: Jagadeesh Dharmudu | Hackathon Event Winning Blueprint"
    p4.font.size = Pt(14)
    p4.font.color.rgb = GREEN_COLOR
    p4.font.name = "Segoe UI"

    # -----------------------------------------------------------------
    # SLIDE 2: Functional Architecture Diagram (Light Theme with Shapes)
    # -----------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_BG
    bg2.line.fill.background()

    # Header
    hdrBox = slide2.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
    htf = hdrBox.text_frame
    hp = htf.paragraphs[0]
    hp.text = "End-to-End Functional Architecture Flow"
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = TEXT_DARK
    hp.font.name = "Segoe UI"

    # 5 Phase Boxes
    phases = [
        {"num": "1", "title": "DETECTION & INGEST", "desc": "Defender for Cloud\nLog Analytics KQL\nAzure Monitor Metrics\n\nTelemetry triggers Event Grid Webhook.", "color": RGBColor(254, 244, 255)},
        {"num": "2", "title": "COGNITIVE SWARM", "desc": "Azure Function V2\nAzure OpenAI (GPT)\n\nExtracts incident, locates file, patches HCL and validates syntax.", "color": RGBColor(239, 246, 255)},
        {"num": "3", "title": "GIT CONTROL PLANE", "desc": "Target Repository\nSecurity Branch\nPull Request Gateway\n\nCreates a human-in-the-loop SRE gate.", "color": RGBColor(240, 253, 244)},
        {"num": "4", "title": "VERIFY & DEPLOY", "desc": "CI/CD Pipeline\nterraform plan\nterraform apply\n\nValidation check triggers on PR merge approval.", "color": RGBColor(254, 242, 242)},
        {"num": "5", "title": "LIVE REMEDIATION", "desc": "Azure Resource Group\nLanding Zone Assets\n\nLive state synced to Git source of truth.", "color": RGBColor(254, 252, 232)}
    ]

    left_margin = Inches(0.8)
    width = Inches(2.1)
    height = Inches(4.5)
    gap = Inches(0.3)
    top_margin = Inches(1.8)

    for i, phase in enumerate(phases):
        x = left_margin + i * (width + gap)
        # Create shape box
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_margin, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = phase["color"]
        box.line.color.rgb = TEXT_MUTED
        box.line.width = Pt(1.5)

        # Add text
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.TOP
        btf.word_wrap = True
        
        bp1 = btf.paragraphs[0]
        bp1.text = f"PHASE {phase['num']}"
        bp1.font.size = Pt(12)
        bp1.font.bold = True
        bp1.font.color.rgb = TEXT_MUTED
        bp1.alignment = PP_ALIGN.CENTER
        bp1.space_after = Pt(10)

        bp2 = btf.add_paragraph()
        bp2.text = phase["title"]
        bp2.font.size = Pt(14)
        bp2.font.bold = True
        bp2.font.color.rgb = TEXT_DARK
        bp2.alignment = PP_ALIGN.CENTER
        bp2.space_after = Pt(15)

        bp3 = btf.add_paragraph()
        bp3.text = phase["desc"]
        bp3.font.size = Pt(11)
        bp3.font.color.rgb = TEXT_DARK
        bp3.alignment = PP_ALIGN.LEFT
        
        # Connectors (except for the last box)
        if i < 4:
            arrow = slide2.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, 
                x + width + Inches(0.05), 
                top_margin + Inches(2.0), 
                Inches(0.2), 
                Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL_COLOR
            arrow.line.fill.background()

    # -----------------------------------------------------------------
    # SLIDE 3: Deep Dive Mappings & Specifications
    # -----------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = LIGHT_BG
    bg3.line.fill.background()

    # Header
    hdrBox3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
    htf3 = hdrBox3.text_frame
    hp3 = htf3.paragraphs[0]
    hp3.text = "Deep Dive Lifecycle Mapping & Specifications"
    hp3.font.size = Pt(28)
    hp3.font.bold = True
    hp3.font.color.rgb = TEXT_DARK
    hp3.font.name = "Segoe UI"

    # Add Table
    rows = 5
    cols = 3
    left = Inches(1.0)
    top = Inches(1.6)
    width = Inches(11.333)
    height = Inches(4.8)

    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(2.333)
    table.columns[1].width = Inches(5.0)
    table.columns[2].width = Inches(4.0)

    # Table Header
    headers = ["Lifecycle Phase", "Technical Interface & Mechanics", "OpenAI Cognitive Footprint"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_LIGHT
            paragraph.font.name = "Segoe UI"

    # Table Content
    data = [
        [
            "DIAGNOSING & INGESTION",
            "Event Grid standardizes polymorphic JSON structures from Defender for Cloud & Monitor via the CloudEvents v1.0 standard, routing it to secure function endpoints.",
            "Footprint #1: Cognitive Parser\nExtracts variables and locates target Landing Zone files dynamically."
        ],
        [
            "REMEDIATING & ORCHESTRATION",
            "Azure Function V2 checks out repo, passes the HCL file to OpenAI, commits the patched declarative code and programmatically spawns a Pull Request.",
            "Footprint #2: Generative HCL Patching\nSurgically updates compliance settings (e.g. disabling public network security ports)."
        ],
        [
            "VERIFYING & CI/CD",
            "Triggered dynamically on PR creation. Runs terraform plan to verify HCL syntax validity and show state blast radius before deployment.",
            "None\n(Decoupled completely from LLM generation to guarantee safety)."
        ],
        [
            "REPORTING & AUDIT",
            "Merges Git branch upon SRE manual approval. Generates an immutable, cryptographically signed ledger record inside the Git history.",
            "Footprint #3: PR Copywriting\nDrafts technical markdown risk analysis and validation guidance inside the SRE PR."
        ]
    ]

    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else ACCENT_BOX
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_DARK
                paragraph.font.name = "Segoe UI"
                if col_idx == 0:
                    paragraph.font.bold = True

    # -----------------------------------------------------------------
    # SLIDE 4: Hackathon Pitch & Value Proposition (Light Theme)
    # -----------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = LIGHT_BG
    bg4.line.fill.background()

    # Header
    hdrBox4 = slide4.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
    htf4 = hdrBox4.text_frame
    hp4 = htf4.paragraphs[0]
    hp4.text = "Hackathon Winning Pitch & Core Value Props"
    hp4.font.size = Pt(28)
    hp4.font.bold = True
    hp4.font.color.rgb = TEXT_DARK
    hp4.font.name = "Segoe UI"

    # Left Column: MTTR Stats
    leftBox = slide4.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
    ltf = leftBox.text_frame
    ltf.word_wrap = True

    lp1 = ltf.paragraphs[0]
    lp1.text = "THE METRIC THAT WINS: ZERO-MTTR"
    lp1.font.size = Pt(16)
    lp1.font.bold = True
    lp1.font.color.rgb = TEAL_COLOR
    lp1.space_after = Pt(15)

    lp2 = ltf.add_paragraph()
    lp2.text = "Traditional Security Remediation: 24 to 72 Hours"
    lp2.font.size = Pt(14)
    lp2.font.bold = True
    lp2.font.color.rgb = TEXT_DARK
    lp2.space_after = Pt(5)

    lp3 = ltf.add_paragraph()
    lp3.text = "Requires ticketing system, human scheduling, manual discovery, and manual command line intervention."
    lp3.font.size = Pt(12)
    lp3.font.color.rgb = TEXT_MUTED
    lp3.space_after = Pt(20)

    lp4 = ltf.add_paragraph()
    lp4.text = "SecOps Swarm Remediation: Under 10 Minutes total"
    lp4.font.size = Pt(14)
    lp4.font.bold = True
    lp4.font.color.rgb = GREEN_COLOR
    lp4.space_after = Pt(5)

    lp5 = ltf.add_paragraph()
    lp5.text = "• Event Ingest & Routing: 1 - 3 Minutes\n• AI Triage & PR Generation: 20 Seconds\n• SRE Human Approval: 2 - 5 Minutes\n• CI/CD Execution: 1 - 2 Minutes"
    lp5.font.size = Pt(12)
    lp5.font.color.rgb = TEXT_DARK
    lp5.space_after = Pt(20)

    # Right Column: Value Propositions
    rightBox = slide4.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.0))
    rtf = rightBox.text_frame
    rtf.word_wrap = True

    rp1 = rtf.paragraphs[0]
    rp1.text = "CORE ARCHITECTURAL PILLARS"
    rp1.font.size = Pt(16)
    rp1.font.bold = True
    rp1.font.color.rgb = TEAL_COLOR
    rp1.space_after = Pt(15)

    pillars = [
        ("1. Shift-Left & Shield-Right", "We never modify live cloud environments directly. We patch the Infrastructure as Code (IaC) repository. Git remains the Single Source of Truth (SSOT)."),
        ("2. Human-in-the-Loop Safeguards", "Instead of direct automation, we open a Pull Request. SREs verify the change in standard workflows before it is merged and deployed."),
        ("3. Dynamic Landing Zone Targeting", "AI scans the repository directory tree dynamically to map alerts to Storage, Network, Compute, or Database files, covering all Azure landing zone structures.")
    ]

    for title, desc in pillars:
        r_p = rtf.add_paragraph()
        r_p.text = title
        r_p.font.size = Pt(14)
        r_p.font.bold = True
        r_p.font.color.rgb = TEXT_DARK
        r_p.space_after = Pt(3)

        r_p2 = rtf.add_paragraph()
        r_p2.text = desc
        r_p2.font.size = Pt(11)
        r_p2.font.color.rgb = TEXT_MUTED
        r_p2.space_after = Pt(15)

    # Save Presentation
    prs.save("references/azureops_gitops_architecture.pptx")
    print("[*] Presentation generated successfully: references/azureops_gitops_architecture.pptx")

if __name__ == "__main__":
    create_presentation()
